from flask import Flask, render_template, request, jsonify, redirect, url_for, session, send_from_directory, make_response
from flask_login import LoginManager, login_user, login_required, logout_user, current_user
from werkzeug.security import generate_password_hash, check_password_hash
from datetime import datetime
import requests
import os
import shutil
import sqlite3
import glob
import secrets
import hmac
import time
import threading
import hashlib
import json
import io
from collections import defaultdict, deque
from urllib.parse import quote
import markdown
import bleach
from markdown.extensions.codehilite import CodeHiliteExtension
from markdown.extensions.fenced_code import FencedCodeExtension
from markdown.extensions.tables import TableExtension
from markdown.extensions.toc import TocExtension

# 导入数据库模型
from database import db, User, ChatHistory, Attachment, RuleDocument, RuleReviewMessage, format_datetime_value

# 初始化 Flask
app = Flask(__name__)

# ============ 配置部分 ============
def env_flag(name, default=False):
    value = os.environ.get(name)
    if value is None:
        return default
    return value.lower() in ('1', 'true', 'yes', 'on')

base_dir = os.path.abspath(os.path.dirname(__file__))
default_data_dir = os.path.abspath(os.path.expanduser(os.environ.get('APP_DATA_DIR', os.path.join(base_dir, 'app_data'))))
default_db_path = os.path.join(default_data_dir, 'app.db')
configured_db_path = os.path.expanduser(os.environ.get('APP_DB_PATH', default_db_path))
configured_db_path = os.path.abspath(configured_db_path)
os.makedirs(os.path.dirname(configured_db_path), exist_ok=True)


def resolve_secret_key():
    env_secret = os.environ.get('SECRET_KEY')
    if env_secret:
        return env_secret

    key_file = os.path.abspath(
        os.path.expanduser(
            os.environ.get('APP_SECRET_KEY_FILE', os.path.join(default_data_dir, '.secret_key'))
        )
    )
    os.makedirs(os.path.dirname(key_file), exist_ok=True)

    if os.path.exists(key_file):
        try:
            with open(key_file, 'r', encoding='utf-8') as f:
                existing = f.read().strip()
            if existing:
                return existing
        except Exception:
            app.logger.exception('读取 SECRET_KEY 文件失败: %s', key_file)

    generated = secrets.token_urlsafe(64)
    try:
        with open(key_file, 'x', encoding='utf-8') as f:
            f.write(generated)
        os.chmod(key_file, 0o600)
        print(f"⚠️ 未设置 SECRET_KEY，已生成并保存到: {key_file}")
        return generated
    except FileExistsError:
        try:
            with open(key_file, 'r', encoding='utf-8') as f:
                existing = f.read().strip()
            if existing:
                return existing
        except Exception:
            app.logger.exception('读取已存在的 SECRET_KEY 文件失败: %s', key_file)
    except Exception:
        app.logger.exception('写入 SECRET_KEY 文件失败: %s', key_file)

    print("⚠️ 未设置 SECRET_KEY，已使用进程内随机密钥（重启后可能失效）")
    return generated


secret_key = resolve_secret_key()

def _table_count(db_file, table_name):
    if not os.path.exists(db_file):
        return 0
    conn = None
    try:
        conn = sqlite3.connect(db_file)
        cursor = conn.cursor()
        cursor.execute(
            "SELECT 1 FROM sqlite_master WHERE type='table' AND name=?",
            (table_name,)
        )
        if cursor.fetchone() is None:
            return 0
        cursor.execute(f"SELECT COUNT(1) FROM {table_name}")
        row = cursor.fetchone()
        return int(row[0]) if row else 0
    except Exception:
        return 0
    finally:
        if conn is not None:
            conn.close()


def _db_table_stats(db_file):
    users_count = max(_table_count(db_file, 'users'), _table_count(db_file, 'user'))
    chats_count = _table_count(db_file, 'chat_history')
    return users_count, chats_count


def _db_data_score(db_file):
    users_count, chats_count = _db_table_stats(db_file)
    return users_count + chats_count


def _legacy_db_candidates():
    candidates = [
        os.path.join(base_dir, 'app.db'),
        os.path.join(base_dir, 'instance', 'app.db')
    ]
    backup_candidates = sorted(
        glob.glob(os.path.join(base_dir, 'backups', 'app_*.db')),
        reverse=True
    )
    candidates.extend(backup_candidates)
    unique = []
    seen = set()
    for item in candidates:
        normalized = os.path.abspath(item)
        if normalized in seen:
            continue
        seen.add(normalized)
        unique.append(normalized)
    return unique


legacy_candidates = _legacy_db_candidates()

best_legacy_db = None
best_legacy_users = 0
best_legacy_chats = 0
best_legacy_score = 0
for legacy_db in legacy_candidates:
    if legacy_db == configured_db_path:
        continue
    users_count, chats_count = _db_table_stats(legacy_db)
    score = users_count + chats_count
    if score > best_legacy_score:
        best_legacy_db = legacy_db
        best_legacy_users = users_count
        best_legacy_chats = chats_count
        best_legacy_score = score

# 迁移策略：
# 1) 目标库不存在时，直接复制首个存在且有数据的旧库。
# 2) 目标库已存在但数据更少（尤其聊天记录更少）时，尝试从更完整旧库补迁。
if best_legacy_db and best_legacy_score > 0:
    if not os.path.exists(configured_db_path):
        shutil.copy2(best_legacy_db, configured_db_path)
        print(f"✓ 已迁移数据库到持久目录: {configured_db_path}")
    else:
        current_users, current_chats = _db_table_stats(configured_db_path)
        current_score = current_users + current_chats
        should_replace = (
            current_score == 0 or
            (best_legacy_chats > current_chats and best_legacy_score > current_score)
        )
        if should_replace:
            backup_file = f"{configured_db_path}.pre-legacy-sync.bak"
            if not os.path.exists(backup_file):
                shutil.copy2(configured_db_path, backup_file)
            shutil.copy2(best_legacy_db, configured_db_path)
            print(
                "✓ 已从旧库补迁历史数据: "
                f"{best_legacy_db} -> {configured_db_path} "
                f"(legacy users/chats={best_legacy_users}/{best_legacy_chats}, "
                f"current users/chats={current_users}/{current_chats})"
            )

app.config['SECRET_KEY'] = secret_key
app.config['SQLALCHEMY_DATABASE_URI'] = f"sqlite:///{configured_db_path}"
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['OLLAMA_BASE_URL'] = os.environ.get('OLLAMA_BASE_URL', 'http://localhost:11434')
app.config['DEFAULT_MODEL'] = os.environ.get('DEFAULT_MODEL', '').strip()
app.config['MAX_CONTENT_LENGTH'] = int(os.environ.get('MAX_CONTENT_LENGTH', '26214400'))  # 25MB
app.config['MAX_QUESTION_CHARS'] = int(os.environ.get('MAX_QUESTION_CHARS', '8000'))
app.config['MAX_PROMPT_CHARS'] = int(os.environ.get('MAX_PROMPT_CHARS', '32000'))
app.config['MAX_ATTACHMENT_SIZE_BYTES'] = int(os.environ.get('MAX_ATTACHMENT_SIZE_BYTES', '20971520'))  # 20MB
app.config['MAX_ATTACHMENTS_PER_REQUEST'] = int(os.environ.get('MAX_ATTACHMENTS_PER_REQUEST', '5'))
app.config['MAX_ATTACHMENTS_PER_MESSAGE'] = int(os.environ.get('MAX_ATTACHMENTS_PER_MESSAGE', '5'))
app.config['MAX_ATTACHMENT_TEXT_CHARS'] = int(os.environ.get('MAX_ATTACHMENT_TEXT_CHARS', '12000'))
app.config['MAX_RULE_FILE_SIZE_BYTES'] = int(os.environ.get('MAX_RULE_FILE_SIZE_BYTES', '10485760'))  # 10MB
app.config['MAX_RULE_FILES_PER_REQUEST'] = int(os.environ.get('MAX_RULE_FILES_PER_REQUEST', '3'))
app.config['MAX_RULE_TEXT_CHARS'] = int(os.environ.get('MAX_RULE_TEXT_CHARS', '20000'))
app.config['SESSION_COOKIE_HTTPONLY'] = True
app.config['SESSION_COOKIE_SAMESITE'] = os.environ.get('SESSION_COOKIE_SAMESITE', 'Lax')
app.config['SESSION_COOKIE_SECURE'] = env_flag('SESSION_COOKIE_SECURE', False)
app.config['REMEMBER_COOKIE_HTTPONLY'] = True
app.config['REMEMBER_COOKIE_SAMESITE'] = os.environ.get('REMEMBER_COOKIE_SAMESITE', 'Lax')
app.config['REMEMBER_COOKIE_SECURE'] = env_flag('REMEMBER_COOKIE_SECURE', app.config['SESSION_COOKIE_SECURE'])

attachment_base_dir = os.path.join(os.path.dirname(configured_db_path), 'uploads')
os.makedirs(attachment_base_dir, exist_ok=True)


class InMemoryRateLimiter:
    def __init__(self):
        self._events = defaultdict(deque)
        self._lock = threading.Lock()

    def allow(self, key, limit, window_seconds):
        now = time.time()
        cutoff = now - window_seconds
        with self._lock:
            bucket = self._events[key]
            while bucket and bucket[0] <= cutoff:
                bucket.popleft()
            if len(bucket) >= limit:
                retry_after = max(1, int(window_seconds - (now - bucket[0])))
                return False, retry_after
            bucket.append(now)
        return True, 0


rate_limiter = InMemoryRateLimiter()

# ============ 初始化扩展 ============
db.init_app(app)
login_manager = LoginManager(app)
login_manager.login_view = 'login'
login_manager.login_message = '请先登录'


@login_manager.unauthorized_handler
def handle_unauthorized():
    if request.path.startswith('/api/'):
        return jsonify({'success': False, 'message': '未登录或会话已失效'}), 401
    return redirect(url_for('login'))

# ============ 用户加载器 ============
@login_manager.user_loader
def load_user(user_id):
    return User.query.get(int(user_id))


def get_csrf_token():
    token = session.get('csrf_token')
    if not token:
        token = secrets.token_urlsafe(32)
        session['csrf_token'] = token
    return token


def get_client_ip():
    forwarded = request.headers.get('X-Forwarded-For', '')
    if forwarded:
        return forwarded.split(',')[0].strip()
    return request.remote_addr or 'unknown'


def enforce_rate_limit(scope, identity, limit, window_seconds=60):
    key = f"{scope}:{identity}"
    allowed, retry_after = rate_limiter.allow(key, limit, window_seconds)
    if allowed:
        return None
    response = jsonify({'success': False, 'message': '请求过于频繁，请稍后再试'})
    response.status_code = 429
    response.headers['Retry-After'] = str(retry_after)
    return response


def get_ollama_base_urls():
    configured = (app.config.get('OLLAMA_BASE_URL') or '').strip()
    candidates = []
    for url in [configured, 'http://127.0.0.1:11434', 'http://localhost:11434']:
        if not url:
            continue
        normalized = url.rstrip('/')
        if normalized not in candidates:
            candidates.append(normalized)
    return candidates


def extract_model_names(payload):
    if not isinstance(payload, dict):
        return []
    models = payload.get('models')
    if not isinstance(models, list):
        return []

    model_names = []
    for model in models:
        if isinstance(model, str):
            model_names.append(model)
            continue
        if not isinstance(model, dict):
            continue
        # Ollama 不同版本可能返回 name 或 model 字段
        name = (model.get('name') or model.get('model') or '').strip()
        if name:
            model_names.append(name)
    return model_names


def safe_requests_get(url, timeout):
    # 忽略环境代理，避免 localhost 请求被 http_proxy/https_proxy 劫持
    with requests.Session() as session:
        session.trust_env = False
        return session.get(url, timeout=timeout)


def safe_requests_post(url, payload, timeout):
    # 忽略环境代理，避免 localhost 请求被 http_proxy/https_proxy 劫持
    with requests.Session() as session:
        session.trust_env = False
        return session.post(url, json=payload, timeout=timeout)


ALLOWED_ATTACHMENT_EXTENSIONS = {
    'pdf', 'txt', 'md',
    'docx', 'xlsx', 'pptx'
}
ALLOWED_RULE_EXTENSIONS = {
    'pdf', 'txt', 'md',
    'docx', 'xlsx', 'pptx'
}


def normalize_extension(filename):
    if not filename or '.' not in filename:
        return ''
    ext = filename.rsplit('.', 1)[-1].strip().lower()
    return ext[:16]


def is_attachment_signature_valid(ext, file_bytes):
    if not file_bytes:
        return False
    signatures = {
        'pdf': [b'%PDF-'],
        'docx': [b'PK\x03\x04', b'PK\x05\x06', b'PK\x07\x08'],
        'xlsx': [b'PK\x03\x04', b'PK\x05\x06', b'PK\x07\x08'],
        'pptx': [b'PK\x03\x04', b'PK\x05\x06', b'PK\x07\x08']
    }
    if ext in signatures:
        return any(file_bytes.startswith(sig) for sig in signatures[ext])
    return True


def decode_text_bytes(file_bytes):
    for encoding in ('utf-8-sig', 'utf-8', 'gb18030', 'big5', 'latin-1'):
        try:
            return file_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode('utf-8', errors='ignore')


def extract_attachment_text(ext, file_bytes):
    try:
        if ext in {'txt', 'md'}:
            return decode_text_bytes(file_bytes), ''

        if ext == 'pdf':
            try:
                from pypdf import PdfReader
            except Exception:
                return '', '缺少 pypdf 依赖'
            reader = PdfReader(io.BytesIO(file_bytes))
            pages = []
            for page in reader.pages:
                pages.append(page.extract_text() or '')
            return '\n'.join(pages), ''

        if ext == 'docx':
            try:
                from docx import Document
            except Exception:
                return '', '缺少 python-docx 依赖'
            doc = Document(io.BytesIO(file_bytes))
            paragraphs = [p.text for p in doc.paragraphs if (p.text or '').strip()]
            table_text = []
            for table in doc.tables:
                for row in table.rows:
                    cells = [(cell.text or '').strip() for cell in row.cells]
                    if any(cells):
                        table_text.append(' | '.join(cells))
            merged = paragraphs + table_text
            return '\n'.join(merged), ''

        if ext == 'xlsx':
            try:
                import openpyxl
            except Exception:
                return '', '缺少 openpyxl 依赖'
            workbook = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True, read_only=True)
            lines = []
            for sheet in workbook.worksheets:
                lines.append(f"## Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    values = [str(v).strip() for v in row if v is not None and str(v).strip()]
                    if values:
                        lines.append(' | '.join(values))
            return '\n'.join(lines), ''

        if ext == 'pptx':
            try:
                from pptx import Presentation
            except Exception:
                return '', '缺少 python-pptx 依赖'
            presentation = Presentation(io.BytesIO(file_bytes))
            lines = []
            for idx, slide in enumerate(presentation.slides, start=1):
                lines.append(f"## Slide {idx}")
                for shape in slide.shapes:
                    text = getattr(shape, 'text', '')
                    if (text or '').strip():
                        lines.append(text.strip())
            return '\n'.join(lines), ''

        return '', '不支持的文件类型'
    except Exception as exc:
        return '', str(exc)


def normalize_text_for_markdown(text):
    normalized = (text or '').replace('\r\n', '\n').replace('\r', '\n')
    lines = []
    blank_run = 0
    for raw in normalized.split('\n'):
        line = raw.rstrip()
        if not line.strip():
            blank_run += 1
            if blank_run > 1:
                continue
            lines.append('')
            continue
        blank_run = 0
        lines.append(line)
    return '\n'.join(lines).strip()


def convert_rule_source_to_markdown(file_name, ext, extracted_text, max_chars):
    safe_name = (file_name or '').replace('\x00', '').replace('\r', ' ').replace('\n', ' ').strip()[:255]
    if not safe_name:
        safe_name = f'rule.{ext or "txt"}'
    safe_ext = (ext or normalize_extension(safe_name) or 'txt').strip().lower()[:16]
    normalized = normalize_text_for_markdown(extracted_text)
    if not normalized:
        return ''

    lines = normalized.split('\n')
    transformed = []
    for line in lines:
        stripped = line.strip()
        if safe_ext == 'xlsx' and stripped.startswith('## Sheet:'):
            sheet_title = stripped.split(':', 1)[1].strip() if ':' in stripped else ''
            transformed.append(f"## 工作表: {sheet_title or '未命名'}")
            continue
        if safe_ext == 'pptx' and stripped.startswith('## Slide '):
            transformed.append(stripped.replace('## Slide ', '## 幻灯片 ', 1))
            continue
        transformed.append(line)
    body = '\n'.join(transformed).strip()

    if safe_ext == 'md':
        markdown_text = body
    else:
        title = os.path.splitext(safe_name)[0].strip() or '规则文档'
        markdown_text = (
            f"# 规则文档：{title}\n\n"
            f"- 来源文件：`{safe_name}`\n"
            f"- 来源格式：`{safe_ext}`\n"
            "- 规范格式：`markdown`\n\n"
            "## 规则正文\n\n"
            f"{body}\n"
        )

    limit = max(2000, int(max_chars or 0))
    if len(markdown_text) > limit:
        markdown_text = markdown_text[:limit] + '\n...(内容已截断)'
    return markdown_text.strip()


def get_rule_markdown_content(rule):
    raw = (getattr(rule, 'content_text', '') or '').strip()
    if not raw:
        return ''
    max_len = max(2000, app.config['MAX_RULE_TEXT_CHARS'])
    content_format = (getattr(rule, 'content_format', '') or 'plain').strip().lower()
    if content_format == 'markdown':
        if len(raw) > max_len:
            return raw[:max_len] + '\n...(内容已截断)'
        return raw
    return convert_rule_source_to_markdown(
        getattr(rule, 'name', '') or 'rule.txt',
        getattr(rule, 'extension', '') or 'txt',
        raw,
        max_len
    )


def parse_attachment_ids(raw_value):
    if isinstance(raw_value, list):
        values = raw_value
    elif isinstance(raw_value, str):
        text = raw_value.strip()
        if not text:
            return []
        try:
            values = json.loads(text)
        except Exception:
            return []
    else:
        return []

    normalized = []
    for value in values:
        try:
            attachment_id = int(value)
        except (TypeError, ValueError):
            continue
        if attachment_id > 0 and attachment_id not in normalized:
            normalized.append(attachment_id)
    return normalized


def attachment_prompt_snippets(attachments, max_total_chars):
    snippets = []
    remaining = max(0, int(max_total_chars))
    for attachment in attachments:
        if remaining <= 0:
            break
        text = (attachment.extracted_text or '').strip()
        if not text:
            continue
        reserve_for_header = 64
        allowed = max(0, remaining - reserve_for_header)
        if allowed <= 0:
            break
        chunk = text[:allowed]
        if len(text) > len(chunk):
            chunk += '\n...(内容已截断)'
        header = f"[附件:{attachment.original_name}]"
        snippet = f"{header}\n{chunk}"
        snippets.append(snippet)
        remaining -= len(snippet) + 2
    return snippets


def serialize_download_response(content, format_type, filename):
    if format_type == 'json':
        mime_type = 'application/json; charset=utf-8'
        fallback = 'download.json'
    elif format_type == 'txt':
        mime_type = 'text/plain; charset=utf-8'
        fallback = 'download.txt'
    else:
        mime_type = 'text/markdown; charset=utf-8'
        fallback = 'download.md'

    raw_name = os.path.basename(str(filename or '').strip()).replace('\x00', '')
    raw_name = raw_name.replace('\r', '').replace('\n', '').strip()[:180]
    if not raw_name:
        raw_name = fallback

    safe_ascii_name = []
    for ch in raw_name:
        if ('a' <= ch.lower() <= 'z') or ('0' <= ch <= '9') or ch in {'-', '_', '.'}:
            safe_ascii_name.append(ch)
        else:
            safe_ascii_name.append('_')
    safe_ascii_name = ''.join(safe_ascii_name).strip('._')
    if not safe_ascii_name:
        safe_ascii_name = fallback

    response = make_response(content)
    response.headers['Content-Type'] = mime_type
    response.headers['Content-Disposition'] = (
        f'attachment; filename="{safe_ascii_name}"; '
        f"filename*=UTF-8''{quote(raw_name, safe='')}"
    )
    return response


def extract_json_object(text):
    if not text:
        return None
    start = text.find('{')
    end = text.rfind('}')
    if start < 0 or end <= start:
        return None
    candidate = text[start:end + 1]
    try:
        return json.loads(candidate)
    except Exception:
        return None


def build_rule_review_prompt(rule_name, rule_text):
    snippet = (rule_text or '')[:app.config['MAX_RULE_TEXT_CHARS']]
    return (
        "你是规则文件审查助手。请检查规则是否自相矛盾、歧义、不可执行、"
        "缺少关键边界条件、存在安全/合规风险。\n"
        "必须只输出 JSON，不要输出任何额外解释。\n"
        "JSON结构如下：\n"
        "{\n"
        '  "pass": true/false,\n'
        '  "summary": "一句话总结",\n'
        '  "issues": [\n'
        '    {"severity":"high|medium|low","title":"问题标题","detail":"问题说明"}\n'
        "  ],\n"
        '  "revision_suggestion": "建议修改后的规则文本（可选）"\n'
        "}\n\n"
        f"规则文件名: {rule_name}\n"
        "规则正文如下：\n"
        f"{snippet}"
    )


def normalize_rule_review_result(raw_text):
    payload = extract_json_object(raw_text or '')
    if not isinstance(payload, dict):
        return {
            'pass': False,
            'summary': 'AI 返回格式不符合预期，请手工检查',
            'issues': [{'severity': 'medium', 'title': '输出格式错误', 'detail': (raw_text or '')[:300]}],
            'revision_suggestion': ''
        }

    issues = payload.get('issues')
    normalized_issues = []
    if isinstance(issues, list):
        for item in issues[:20]:
            if not isinstance(item, dict):
                continue
            severity = (item.get('severity') or 'medium').strip().lower()
            if severity not in {'high', 'medium', 'low'}:
                severity = 'medium'
            normalized_issues.append({
                'severity': severity,
                'title': str(item.get('title') or '未命名问题')[:80],
                'detail': str(item.get('detail') or '')[:600]
            })

    summary = str(payload.get('summary') or '').strip()
    if not summary:
        summary = '未提供总结'

    return {
        'pass': bool(payload.get('pass')),
        'summary': summary[:300],
        'issues': normalized_issues,
        'revision_suggestion': str(payload.get('revision_suggestion') or '')[:4000]
    }


def get_rule_for_current_user(rule_id, require_current=True):
    query = RuleDocument.query.filter_by(id=rule_id, user_id=current_user.id)
    if require_current:
        query = query.filter_by(is_current=True)
    return query.first()


def get_rule_review_transcript(rule_id, max_rounds=16):
    messages = (
        RuleReviewMessage.query
        .filter_by(user_id=current_user.id, rule_id=rule_id)
        .order_by(RuleReviewMessage.created_at.asc(), RuleReviewMessage.id.asc())
        .all()
    )
    if max_rounds and len(messages) > max_rounds:
        messages = messages[-max_rounds:]
    return messages


def build_rule_review_chat_prompt(rule, transcript, user_message):
    rule_text = get_rule_markdown_content(rule)
    lines = []
    for msg in transcript[-12:]:
        role = '用户' if msg.role == 'user' else '审查助手'
        content = (msg.content or '').strip()[:1200]
        if content:
            lines.append(f"{role}: {content}")
    history_text = '\n'.join(lines) if lines else '（无历史对话）'
    return (
        "你是“规则审查助手”。目标：帮助用户把规则文档修改到可执行、无明显冲突、边界清晰。\n"
        "规则正文统一为 Markdown；你给出的修订稿也必须是 Markdown。\n"
        "请使用自然语言回复，不要求 JSON。\n"
        "输出风格：\n"
        "1) 先给当前结论（通过/未通过）\n"
        "2) 列出关键问题（若有）\n"
        "3) 给出可直接修改的建议条款\n"
        "4) 若已接近通过，明确告诉用户还差什么\n\n"
        "如果用户明确要求“直接修改文档/给出最终稿”，请额外给出“修订稿全文”，"
        "并用 ```markdown ... ``` 包裹，内容可直接保存为文件。\n\n"
        f"规则名称: {rule.name}\n"
        "规则正文:\n"
        f"{rule_text}\n\n"
        "历史对话:\n"
        f"{history_text}\n\n"
        "用户本轮问题:\n"
        f"{user_message}\n"
    )


def build_rule_review_verdict_prompt(rule, transcript):
    rule_text = get_rule_markdown_content(rule)
    lines = []
    for msg in transcript[-20:]:
        role = '用户' if msg.role == 'user' else '审查助手'
        content = (msg.content or '').strip()[:1200]
        if content:
            lines.append(f"{role}: {content}")
    history_text = '\n'.join(lines) if lines else '（无历史对话）'
    return (
        "你是规则审查最终判定器。\n"
        "请基于规则正文和审核对话，判断该规则是否可以进入“确认通过”。\n"
        "规则正文为 Markdown；如给出 revision_suggestion，必须返回 Markdown。\n"
        "必须仅输出 JSON，不要输出其他内容。\n"
        "JSON结构：\n"
        "{\n"
        '  "pass": true/false,\n'
        '  "summary": "一句话总结",\n'
        '  "issues": [\n'
        '    {"severity":"high|medium|low","title":"问题标题","detail":"问题说明"}\n'
        "  ],\n"
        '  "revision_suggestion": "建议修改后的规则文本（可选）"\n'
        "}\n\n"
        f"规则名称: {rule.name}\n"
        "规则正文:\n"
        f"{rule_text}\n\n"
        "审核对话:\n"
        f"{history_text}\n"
    )


def format_rule_verdict_message(verdict):
    result_label = '通过' if verdict.get('pass') else '未通过'
    summary = (verdict.get('summary') or '').strip() or '无总结'
    issues = verdict.get('issues') or []
    rows = [f"审核结论：{result_label}", f"总结：{summary}"]
    if issues:
        rows.append("主要问题：")
        for item in issues[:6]:
            sev = item.get('severity') or 'medium'
            title = item.get('title') or '未命名问题'
            detail = item.get('detail') or ''
            rows.append(f"- [{sev}] {title}：{detail}")
    suggestion = (verdict.get('revision_suggestion') or '').strip()
    if suggestion:
        rows.append("建议修改稿：")
        rows.append(suggestion[:2000])
    return '\n'.join(rows)


@app.before_request
def csrf_protect():
    if request.method not in {'POST', 'PUT', 'PATCH', 'DELETE'}:
        return None

    # 仅对已登录用户强制校验，避免破坏登录/注册流程
    if not current_user.is_authenticated:
        return None

    provided_token = request.headers.get('X-CSRF-Token') or request.form.get('csrf_token')
    expected_token = session.get('csrf_token')
    if not expected_token or not provided_token or not hmac.compare_digest(provided_token, expected_token):
        return jsonify({'success': False, 'message': 'CSRF token 校验失败'}), 403
    return None


@app.after_request
def add_security_headers(response):
    response.headers['X-Content-Type-Options'] = 'nosniff'
    response.headers['X-Frame-Options'] = 'DENY'
    response.headers['Referrer-Policy'] = 'no-referrer'
    return response


@app.errorhandler(413)
def payload_too_large(_):
    return jsonify({'success': False, 'message': '请求体过大'}), 413

# ============ Markdown渲染函数 ============
def render_markdown(text):
    """将Markdown文本渲染为HTML"""
    if not text:
        return ''
    
    # 配置Markdown扩展
    extensions = [
        'markdown.extensions.extra',
        'markdown.extensions.abbr',
        'markdown.extensions.attr_list',
        'markdown.extensions.def_list',
        'markdown.extensions.footnotes',
        'markdown.extensions.tables',
        'markdown.extensions.admonition',
        CodeHiliteExtension(
            css_class='highlight',
            linenums=False,
            guess_lang=True,
            use_pygments=True
        ),
        FencedCodeExtension(),
        TableExtension(),
        TocExtension(toc_depth='2-6')
    ]
    
    # 转换Markdown为HTML
    html = markdown.markdown(text, extensions=extensions)
    
    # 允许的HTML标签和属性
    allowed_tags = [
        'h1', 'h2', 'h3', 'h4', 'h5', 'h6',
        'b', 'strong', 'i', 'em', 'u', 'del',
        'p', 'br', 'hr',
        'ul', 'ol', 'li',
        'blockquote',
        'pre', 'code',
        'table', 'thead', 'tbody', 'tr', 'th', 'td',
        'a', 'img',
        'div', 'span'
    ]
    
    allowed_attrs = {
        '*': ['class', 'id'],
        'a': ['href', 'title', 'target'],
        'img': ['src', 'alt', 'title'],
        'code': ['class'],
        'pre': ['class']
    }
    
    # 清理HTML，防止XSS攻击
    clean_html = bleach.clean(
        html,
        tags=allowed_tags,
        attributes=allowed_attrs,
        strip=True
    )
    
    return clean_html

# ============ 数据库初始化 ============
def ensure_db_schema_compatibility():
    if not os.path.exists(configured_db_path):
        return

    conn = None
    try:
        conn = sqlite3.connect(configured_db_path)
        cursor = conn.cursor()

        def table_columns(table_name):
            cursor.execute(f"PRAGMA table_info({table_name})")
            return {row[1] for row in cursor.fetchall()}

        def table_exists(table_name):
            cursor.execute(
                "SELECT 1 FROM sqlite_master WHERE type='table' AND name=?",
                (table_name,)
            )
            return cursor.fetchone() is not None

        def ensure_column(table_name, column_name, ddl):
            if not table_exists(table_name):
                return
            existing = table_columns(table_name)
            if column_name in existing:
                return
            cursor.execute(f"ALTER TABLE {table_name} ADD COLUMN {ddl}")
            print(f"✓ 已补齐字段: {table_name}.{column_name}")

        ensure_column('users', 'created_at', 'created_at DATETIME')
        ensure_column('users', 'last_login', 'last_login DATETIME')
        ensure_column('users', 'is_active', 'is_active BOOLEAN DEFAULT 1')

        ensure_column('chat_history', 'answer_html', "answer_html TEXT DEFAULT ''")
        ensure_column('chat_history', 'model', "model VARCHAR(100) DEFAULT ''")
        ensure_column('chat_history', 'tokens_used', 'tokens_used INTEGER DEFAULT 0')
        ensure_column('chat_history', 'created_at', 'created_at DATETIME')
        ensure_column('chat_history', 'conversation_id', "conversation_id VARCHAR(64) DEFAULT ''")
        ensure_column('chat_history', 'conversation_title', "conversation_title VARCHAR(120) DEFAULT ''")
        ensure_column('chat_history', 'question_html', "question_html TEXT DEFAULT ''")
        ensure_column('chat_history', 'attachment_ids', "attachment_ids TEXT DEFAULT '[]'")

        ensure_column('rule_documents', 'rule_group_id', "rule_group_id VARCHAR(32) DEFAULT ''")
        ensure_column('rule_documents', 'version', 'version INTEGER DEFAULT 1')
        ensure_column('rule_documents', 'name', "name VARCHAR(255) DEFAULT ''")
        ensure_column('rule_documents', 'extension', "extension VARCHAR(16) DEFAULT ''")
        ensure_column('rule_documents', 'content_format', "content_format VARCHAR(16) DEFAULT 'plain'")
        ensure_column('rule_documents', 'content_text', "content_text TEXT DEFAULT ''")
        ensure_column('rule_documents', 'status', "status VARCHAR(32) DEFAULT 'draft'")
        ensure_column('rule_documents', 'is_current', 'is_current BOOLEAN DEFAULT 1')
        ensure_column('rule_documents', 'is_active', 'is_active BOOLEAN DEFAULT 0')
        ensure_column('rule_documents', 'ai_review_passed', 'ai_review_passed BOOLEAN DEFAULT 0')
        ensure_column('rule_documents', 'ai_review_summary', "ai_review_summary TEXT DEFAULT ''")
        ensure_column('rule_documents', 'ai_review_raw', "ai_review_raw TEXT DEFAULT ''")
        ensure_column('rule_documents', 'created_at', 'created_at DATETIME')
        ensure_column('rule_documents', 'updated_at', 'updated_at DATETIME')

        ensure_column('rule_review_messages', 'user_id', 'user_id INTEGER')
        ensure_column('rule_review_messages', 'rule_id', 'rule_id INTEGER')
        ensure_column('rule_review_messages', 'role', "role VARCHAR(16) DEFAULT 'user'")
        ensure_column('rule_review_messages', 'content', "content TEXT DEFAULT ''")
        ensure_column('rule_review_messages', 'created_at', 'created_at DATETIME')

        if table_exists('rule_documents'):
            cursor.execute(
                "UPDATE rule_documents SET content_format='plain' "
                "WHERE content_format IS NULL OR TRIM(content_format)=''"
            )

        conn.commit()
    except Exception:
        app.logger.exception('数据库兼容字段补齐失败')
    finally:
        if conn is not None:
            conn.close()


def init_db():
    """初始化数据库"""
    with app.app_context():
        ensure_db_schema_compatibility()
        # 创建所有表
        db.create_all()
        print("✓ 数据库表创建成功")
        
        # 检查是否需要初始化管理员
        admin = User.query.filter_by(username='admin').first()
        if not admin:
            admin_init_password = os.environ.get('ADMIN_INIT_PASSWORD')
            if admin_init_password:
                admin = User(
                    username='admin',
                    created_at=datetime.utcnow(),
                    is_active=True
                )
                admin.set_password(admin_init_password)
                db.session.add(admin)
                db.session.commit()
                print("✓ 管理员用户已创建 (admin)")
            else:
                print("⚠️ 未设置 ADMIN_INIT_PASSWORD，已跳过默认管理员创建")
        else:
            print("✓ 管理员用户已存在")

# 执行数据库初始化
init_db()

# ============ Ollama API调用 ============
def query_ollama(prompt, model=None):
    try:
        model = model or app.config['DEFAULT_MODEL']
        if not model:
            return {"error": "未指定模型，请先在Ollama中安装并选择模型"}
        last_error = None
        for base_url in get_ollama_base_urls():
            url = f"{base_url}/api/generate"
            try:
                response = safe_requests_post(url, {
                    "model": model,
                    "prompt": prompt,
                    "stream": False
                }, timeout=300)
                response.raise_for_status()
                result = response.json()
                return {
                    "response": result.get("response", ""),
                    "eval_count": result.get("eval_count", 0)
                }
            except requests.exceptions.RequestException as exc:
                last_error = exc
                continue
        if last_error:
            app.logger.warning('Ollama generate 调用失败: %s', last_error)
        return {"error": "无法连接到Ollama服务，请确认服务地址与端口"}
    except Exception:
        app.logger.exception('Ollama API调用异常')
        return {"error": "Ollama服务调用失败，请稍后重试"}

# ============ 路由 ============
@app.route('/')
def index():
    if not current_user.is_authenticated:
        return redirect(url_for('login'))
    return render_template('index.html', csrf_token=get_csrf_token())

@app.route('/login', methods=['GET', 'POST'])
def login():
    if current_user.is_authenticated:
        return redirect(url_for('index'))
    
    if request.method == 'GET':
        return render_template('login.html')
    
    try:
        ip = get_client_ip()
        limit_response = enforce_rate_limit('login_ip', ip, limit=12, window_seconds=60)
        if limit_response:
            return limit_response

        data = request.get_json(silent=True) or {}
        username = data.get('username', '').strip()
        password = data.get('password', '').strip()
        
        if not username or not password:
            return jsonify({'success': False, 'message': '用户名和密码不能为空'})
        
        if len(username) > 80 or len(password) > 256:
            return jsonify({'success': False, 'message': '输入长度超出限制'})
        
        limit_response = enforce_rate_limit('login_user', f"{ip}:{username.lower()}", limit=6, window_seconds=60)
        if limit_response:
            return limit_response
        
        user = User.query.filter_by(username=username).first()
        
        if user and user.check_password(password):
            user.last_login = datetime.utcnow()
            db.session.commit()
            login_user(user, remember=True)
            return jsonify({'success': True, 'redirect': url_for('index')})
        else:
            return jsonify({'success': False, 'message': '用户名或密码错误'})
    except Exception:
        app.logger.exception('登录处理异常')
        return jsonify({'success': False, 'message': '登录失败，请稍后重试'})

@app.route('/register', methods=['POST'])
def register():
    try:
        limit_response = enforce_rate_limit('register_ip', get_client_ip(), limit=5, window_seconds=300)
        if limit_response:
            return limit_response

        data = request.get_json(silent=True) or {}
        username = data.get('username', '').strip()
        password = data.get('password', '').strip()
        
        if not username or not password:
            return jsonify({'success': False, 'message': '用户名和密码不能为空'})
        
        if len(username) < 3:
            return jsonify({'success': False, 'message': '用户名至少3个字符'})
        
        if len(password) < 6:
            return jsonify({'success': False, 'message': '密码至少6个字符'})
        
        if len(username) > 80 or len(password) > 256:
            return jsonify({'success': False, 'message': '输入长度超出限制'})
        
        if User.query.filter_by(username=username).first():
            return jsonify({'success': False, 'message': '用户名已存在'})
        
        user = User(username=username, created_at=datetime.utcnow(), is_active=True)
        user.set_password(password)
        db.session.add(user)
        db.session.commit()
        
        return jsonify({'success': True, 'message': '注册成功，请登录'})
    except Exception:
        app.logger.exception('注册处理异常')
        return jsonify({'success': False, 'message': '注册失败，请稍后重试'})

@app.route('/change-password', methods=['POST'])
@login_required
def change_password():
    try:
        limit_response = enforce_rate_limit(
            'change_password',
            f"{current_user.id}:{get_client_ip()}",
            limit=8,
            window_seconds=300
        )
        if limit_response:
            return limit_response

        data = request.get_json(silent=True) or {}
        old_password = data.get('old_password', '').strip()
        new_password = data.get('new_password', '').strip()
        
        if not old_password or not new_password:
            return jsonify({'success': False, 'message': '密码不能为空'})
        
        if len(new_password) < 6:
            return jsonify({'success': False, 'message': '新密码至少6个字符'})
        
        if len(old_password) > 256 or len(new_password) > 256:
            return jsonify({'success': False, 'message': '输入长度超出限制'})
        
        if current_user.check_password(old_password):
            current_user.set_password(new_password)
            db.session.commit()
            return jsonify({'success': True, 'message': '密码修改成功'})
        else:
            return jsonify({'success': False, 'message': '原密码错误'})
    except Exception:
        app.logger.exception('修改密码异常')
        return jsonify({'success': False, 'message': '密码修改失败，请稍后重试'})

@app.route('/logout', methods=['POST'])
@login_required
def logout():
    logout_user()
    session.pop('csrf_token', None)
    return jsonify({'success': True, 'redirect': url_for('login')})

@app.route('/api/user/info')
@login_required
def user_info():
    try:
        user_payload = current_user.to_dict()
    except Exception:
        app.logger.exception('用户信息序列化失败，回退基础字段')
        user_payload = {
            'id': current_user.id,
            'username': current_user.username,
            'created_at': str(getattr(current_user, 'created_at', '') or ''),
            'last_login': str(getattr(current_user, 'last_login', '') or '')
        }
    return jsonify({
        'success': True,
        'user': user_payload
    })


@app.route('/api/rules', methods=['GET'])
@login_required
def list_rules():
    try:
        include_history = request.args.get('history', '0').strip() in {'1', 'true', 'yes'}
        query = RuleDocument.query.filter_by(user_id=current_user.id)
        if not include_history:
            query = query.filter_by(is_current=True)
        rules = query.order_by(RuleDocument.is_active.desc(), RuleDocument.updated_at.desc(), RuleDocument.id.desc()).all()
        items = []
        for rule in rules:
            data = rule.to_dict()
            review_count = RuleReviewMessage.query.filter_by(user_id=current_user.id, rule_id=rule.id).count()
            data['review_message_count'] = review_count
            items.append(data)
        return jsonify({
            'success': True,
            'rules': items
        })
    except Exception:
        app.logger.exception('加载规则列表异常')
        return jsonify({'success': False, 'message': '加载规则失败'}), 500


@app.route('/api/rules/upload', methods=['POST'])
@login_required
def upload_rules():
    try:
        limit_response = enforce_rate_limit(
            'rules_upload',
            f"{current_user.id}:{get_client_ip()}",
            limit=12,
            window_seconds=60
        )
        if limit_response:
            return limit_response

        replace_rule_id = request.form.get('replace_rule_id', '').strip()
        files = request.files.getlist('files')
        if not files:
            one = request.files.get('file')
            if one:
                files = [one]
        if not files:
            return jsonify({'success': False, 'message': '未检测到规则文件'}), 400

        if replace_rule_id and len(files) != 1:
            return jsonify({'success': False, 'message': '修订上传只能上传单个文件'}), 400

        if len(files) > app.config['MAX_RULE_FILES_PER_REQUEST']:
            return jsonify({'success': False, 'message': f"单次最多上传 {app.config['MAX_RULE_FILES_PER_REQUEST']} 个规则文件"}), 400

        replace_rule = None
        if replace_rule_id:
            try:
                replace_rule = RuleDocument.query.filter_by(
                    id=int(replace_rule_id),
                    user_id=current_user.id,
                    is_current=True
                ).first()
            except ValueError:
                replace_rule = None
            if not replace_rule:
                return jsonify({'success': False, 'message': '待修订规则不存在'}), 404

        created = []
        errors = []
        max_size = max(1024, app.config['MAX_RULE_FILE_SIZE_BYTES'])
        max_text_len = max(2000, app.config['MAX_RULE_TEXT_CHARS'])

        for file in files:
            raw_name = (file.filename or '').strip()
            if not raw_name:
                errors.append({'name': 'unknown', 'error': '文件名为空'})
                continue

            ext = normalize_extension(raw_name)
            if ext not in ALLOWED_RULE_EXTENSIONS:
                errors.append({'name': raw_name, 'error': f'不支持的规则文件类型: .{ext or "unknown"}'})
                continue

            file_bytes = file.read(max_size + 1)
            if len(file_bytes) > max_size:
                errors.append({'name': raw_name, 'error': f'规则文件过大，单文件最大 {max_size // (1024 * 1024)}MB'})
                continue
            if not file_bytes:
                errors.append({'name': raw_name, 'error': '空文件不允许上传'})
                continue
            if not is_attachment_signature_valid(ext, file_bytes):
                errors.append({'name': raw_name, 'error': '文件签名校验失败'})
                continue

            extracted_text, parse_error = extract_attachment_text(ext, file_bytes)
            extracted_text = (extracted_text or '').strip()
            if parse_error:
                errors.append({'name': raw_name, 'error': f'解析失败: {parse_error}'})
                continue
            if not extracted_text:
                errors.append({'name': raw_name, 'error': '未提取到可用文本'})
                continue
            if len(extracted_text) > max_text_len:
                extracted_text = extracted_text[:max_text_len] + '\n...(内容已截断)'

            clean_name = raw_name.replace('\x00', '').replace('\r', ' ').replace('\n', ' ').strip()[:255]
            if not clean_name:
                clean_name = f'rule.{ext}'

            rule_markdown = convert_rule_source_to_markdown(clean_name, ext, extracted_text, max_text_len)
            if not rule_markdown:
                errors.append({'name': raw_name, 'error': '规则内容转换为 Markdown 失败'})
                continue

            if replace_rule:
                group_id = replace_rule.rule_group_id
                latest = (
                    RuleDocument.query
                    .filter_by(user_id=current_user.id, rule_group_id=group_id)
                    .order_by(RuleDocument.version.desc())
                    .first()
                )
                next_version = (latest.version + 1) if latest else (replace_rule.version + 1)
                RuleDocument.query.filter_by(user_id=current_user.id, rule_group_id=group_id, is_current=True).update({
                    'is_current': False,
                    'is_active': False
                })
            else:
                group_id = secrets.token_hex(8)
                next_version = 1

            item = RuleDocument(
                user_id=current_user.id,
                rule_group_id=group_id,
                version=next_version,
                name=clean_name,
                extension=ext,
                content_format='markdown',
                content_text=rule_markdown,
                status='draft',
                is_current=True,
                is_active=False,
                ai_review_passed=False,
                ai_review_summary='',
                ai_review_raw=''
            )
            db.session.add(item)
            db.session.flush()
            created.append(item.to_dict())

        db.session.commit()
        if not created:
            return jsonify({'success': False, 'message': '上传失败', 'errors': errors}), 400

        return jsonify({
            'success': True,
            'rules': created,
            'errors': errors
        })
    except Exception:
        app.logger.exception('上传规则文件异常')
        return jsonify({'success': False, 'message': '上传规则失败'}), 500


@app.route('/api/rules/<int:rule_id>/review/messages', methods=['GET'])
@login_required
def list_rule_review_messages(rule_id):
    try:
        rule = get_rule_for_current_user(rule_id, require_current=True)
        if not rule:
            return jsonify({'success': False, 'message': '规则不存在'}), 404
        messages = get_rule_review_transcript(rule.id, max_rounds=0)
        return jsonify({
            'success': True,
            'rule': rule.to_dict(),
            'messages': [item.to_dict() for item in messages]
        })
    except Exception:
        app.logger.exception('加载规则审核消息异常')
        return jsonify({'success': False, 'message': '加载审核消息失败'}), 500


@app.route('/api/rules/<int:rule_id>/review/messages', methods=['POST'])
@login_required
def create_rule_review_message(rule_id):
    try:
        limit_response = enforce_rate_limit(
            'rules_review_chat',
            f"{current_user.id}:{get_client_ip()}",
            limit=20,
            window_seconds=60
        )
        if limit_response:
            return limit_response

        rule = get_rule_for_current_user(rule_id, require_current=True)
        if not rule:
            return jsonify({'success': False, 'message': '规则不存在'}), 404

        data = request.get_json(silent=True) or {}
        model = (data.get('model') or app.config['DEFAULT_MODEL']).strip()
        user_message = (data.get('message') or '').strip()
        if not model:
            return jsonify({'success': False, 'message': '未设置模型，请先选择模型'}), 400
        if len(user_message) > 4000:
            return jsonify({'success': False, 'message': '消息过长，最多4000字符'}), 400

        transcript = get_rule_review_transcript(rule.id, max_rounds=16)
        if not user_message:
            if transcript:
                return jsonify({'success': False, 'message': '消息不能为空'}), 400
            user_message = '请先对这份规则做第一轮审核，指出主要问题并给出修改建议。'

        user_item = RuleReviewMessage(
            user_id=current_user.id,
            rule_id=rule.id,
            role='user',
            content=user_message
        )
        db.session.add(user_item)
        db.session.flush()
        transcript.append(user_item)

        prompt = build_rule_review_chat_prompt(rule, transcript, user_message)
        result = query_ollama(prompt, model)
        if 'error' in result:
            db.session.rollback()
            return jsonify({'success': False, 'message': result['error']}), 500

        assistant_text = (result.get('response') or '').strip() or '未收到有效审查意见，请重试。'
        assistant_item = RuleReviewMessage(
            user_id=current_user.id,
            rule_id=rule.id,
            role='assistant',
            content=assistant_text[:12000]
        )
        db.session.add(assistant_item)
        db.session.commit()

        return jsonify({
            'success': True,
            'rule': rule.to_dict(),
            'messages': [user_item.to_dict(), assistant_item.to_dict()]
        })
    except Exception:
        app.logger.exception('创建规则审核消息异常')
        return jsonify({'success': False, 'message': '审核对话失败'}), 500


@app.route('/api/rules/<int:rule_id>/review/evaluate', methods=['POST'])
@login_required
def evaluate_rule_review(rule_id):
    try:
        limit_response = enforce_rate_limit(
            'rules_review_eval',
            f"{current_user.id}:{get_client_ip()}",
            limit=10,
            window_seconds=60
        )
        if limit_response:
            return limit_response

        rule = get_rule_for_current_user(rule_id, require_current=True)
        if not rule:
            return jsonify({'success': False, 'message': '规则不存在'}), 404

        data = request.get_json(silent=True) or {}
        model = (data.get('model') or app.config['DEFAULT_MODEL']).strip()
        if not model:
            return jsonify({'success': False, 'message': '未设置模型，请先选择模型'}), 400

        transcript = get_rule_review_transcript(rule.id, max_rounds=20)
        verdict_prompt = build_rule_review_verdict_prompt(rule, transcript)
        result = query_ollama(verdict_prompt, model)
        if 'error' in result:
            return jsonify({'success': False, 'message': result['error']}), 500

        raw = result.get('response', '')
        verdict = normalize_rule_review_result(raw)

        rule.ai_review_passed = bool(verdict.get('pass'))
        rule.ai_review_summary = verdict.get('summary') or ''
        rule.ai_review_raw = raw[:12000]
        rule.status = 'ai_review_passed' if rule.ai_review_passed else 'ai_review_failed'
        if not rule.ai_review_passed:
            rule.is_active = False

        assistant_note = RuleReviewMessage(
            user_id=current_user.id,
            rule_id=rule.id,
            role='assistant',
            content=format_rule_verdict_message(verdict)
        )
        db.session.add(assistant_note)
        db.session.commit()

        return jsonify({
            'success': True,
            'rule': rule.to_dict(),
            'verdict': verdict,
            'message': assistant_note.to_dict()
        })
    except Exception:
        app.logger.exception('规则审核判定异常')
        return jsonify({'success': False, 'message': '判定失败'}), 500


@app.route('/api/rules/<int:rule_id>/ai-review', methods=['POST'])
@login_required
def ai_review_rule(rule_id):
    # 兼容旧前端：转发到新版“判定通过”接口
    return evaluate_rule_review(rule_id)


@app.route('/api/rules/<int:rule_id>/confirm', methods=['POST'])
@login_required
def confirm_rule(rule_id):
    try:
        rule = RuleDocument.query.filter_by(id=rule_id, user_id=current_user.id, is_current=True).first()
        if not rule:
            return jsonify({'success': False, 'message': '规则不存在'}), 404
        if not rule.ai_review_passed or rule.status != 'ai_review_passed':
            return jsonify({'success': False, 'message': '规则尚未通过AI审查，无法确认'}), 400

        rule.status = 'confirmed'
        db.session.commit()
        return jsonify({'success': True, 'rule': rule.to_dict()})
    except Exception:
        app.logger.exception('确认规则异常')
        return jsonify({'success': False, 'message': '确认失败'}), 500


@app.route('/api/rules/<int:rule_id>/active', methods=['POST'])
@login_required
def toggle_rule_active(rule_id):
    try:
        rule = RuleDocument.query.filter_by(id=rule_id, user_id=current_user.id, is_current=True).first()
        if not rule:
            return jsonify({'success': False, 'message': '规则不存在'}), 404
        if rule.status != 'confirmed':
            return jsonify({'success': False, 'message': '只有已确认规则才能启用'}), 400

        data = request.get_json(silent=True) or {}
        active = data.get('active')
        if active is None:
            active = not bool(rule.is_active)
        active = bool(active)
        rule.is_active = active
        db.session.commit()
        return jsonify({'success': True, 'rule': rule.to_dict()})
    except Exception:
        app.logger.exception('启停规则异常')
        return jsonify({'success': False, 'message': '启停失败'}), 500


@app.route('/api/attachments', methods=['POST'])
@login_required
def upload_attachments():
    try:
        limit_response = enforce_rate_limit(
            'attachments',
            f"{current_user.id}:{get_client_ip()}",
            limit=20,
            window_seconds=60
        )
        if limit_response:
            return limit_response

        files = request.files.getlist('files')
        if not files:
            single = request.files.get('file')
            if single:
                files = [single]
        if not files:
            return jsonify({'success': False, 'message': '未检测到上传文件'}), 400

        max_count = max(1, app.config['MAX_ATTACHMENTS_PER_REQUEST'])
        if len(files) > max_count:
            return jsonify({'success': False, 'message': f'单次最多上传 {max_count} 个文件'}), 400

        user_dir = os.path.join(attachment_base_dir, str(current_user.id))
        os.makedirs(user_dir, exist_ok=True)

        success_items = []
        failed_items = []
        max_size = max(1024, app.config['MAX_ATTACHMENT_SIZE_BYTES'])
        max_text_len = max(2000, app.config['MAX_ATTACHMENT_TEXT_CHARS'])

        for file in files:
            raw_name = (file.filename or '').strip()
            if not raw_name:
                failed_items.append({'name': 'unknown', 'error': '文件名为空'})
                continue

            ext = normalize_extension(raw_name)
            if ext not in ALLOWED_ATTACHMENT_EXTENSIONS:
                failed_items.append({'name': raw_name, 'error': f'不支持的文件类型: .{ext or "unknown"}'})
                continue

            data = file.read(max_size + 1)
            if len(data) > max_size:
                failed_items.append({'name': raw_name, 'error': f'文件过大，单文件最大 {max_size // (1024 * 1024)}MB'})
                continue
            if not data:
                failed_items.append({'name': raw_name, 'error': '空文件不允许上传'})
                continue
            if not is_attachment_signature_valid(ext, data):
                failed_items.append({'name': raw_name, 'error': '文件签名校验失败'})
                continue

            digest = hashlib.sha256(data).hexdigest()
            safe_name = (raw_name or '').replace('\x00', '').replace('\r', ' ').replace('\n', ' ').strip()
            if not safe_name:
                safe_name = f'file.{ext}'
            random_name = f"{datetime.utcnow().strftime('%Y%m%d%H%M%S')}_{secrets.token_hex(6)}.{ext}"
            stored_path = os.path.join(user_dir, random_name)

            with open(stored_path, 'wb') as f:
                f.write(data)

            extracted_text, parse_error = extract_attachment_text(ext, data)
            parse_status = 'ready' if not parse_error else 'error'
            extracted_text = (extracted_text or '').strip()
            if len(extracted_text) > max_text_len:
                extracted_text = extracted_text[:max_text_len] + '\n...(内容已截断)'

            attachment = Attachment(
                user_id=current_user.id,
                original_name=safe_name[:255],
                stored_name=random_name,
                stored_path=stored_path,
                extension=ext,
                mime_type=file.mimetype or '',
                size_bytes=len(data),
                sha256=digest,
                parse_status=parse_status,
                parse_error=(parse_error or '')[:255],
                extracted_text=extracted_text
            )
            db.session.add(attachment)
            db.session.flush()
            success_items.append(attachment.to_dict())

        db.session.commit()
        if not success_items and failed_items:
            return jsonify({'success': False, 'message': '上传失败', 'errors': failed_items}), 400

        return jsonify({
            'success': True,
            'attachments': success_items,
            'errors': failed_items
        })
    except Exception:
        app.logger.exception('上传附件异常')
        return jsonify({'success': False, 'message': '上传失败，请稍后重试'}), 500


def resolve_chat_attachments(chat):
    attachment_ids = parse_attachment_ids(getattr(chat, 'attachment_ids', '[]'))
    if not attachment_ids:
        return []
    attachments = (
        Attachment.query
        .filter(Attachment.user_id == current_user.id, Attachment.id.in_(attachment_ids))
        .all()
    )
    attachments_by_id = {item.id: item for item in attachments}
    ordered = []
    for attachment_id in attachment_ids:
        item = attachments_by_id.get(attachment_id)
        if item:
            ordered.append(item)
    return ordered


@app.route('/api/rules/<int:rule_id>/review/messages/<int:message_id>/download', methods=['GET'])
@login_required
def download_rule_review_message(rule_id, message_id):
    try:
        format_type = (request.args.get('format') or 'md').strip().lower()
        if format_type not in {'md', 'txt', 'json'}:
            return jsonify({'success': False, 'message': '不支持的导出格式'}), 400

        rule = RuleDocument.query.filter_by(id=rule_id, user_id=current_user.id).first()
        if not rule:
            return jsonify({'success': False, 'message': '规则不存在'}), 404

        message = RuleReviewMessage.query.filter_by(
            id=message_id,
            user_id=current_user.id,
            rule_id=rule.id
        ).first()
        if not message:
            return jsonify({'success': False, 'message': '审核消息不存在'}), 404

        body = (message.content or '').strip()
        if not body:
            return jsonify({'success': False, 'message': '审核消息内容为空'}), 400

        generated_at = datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S')
        base_name = os.path.splitext((rule.name or '').strip())[0]
        safe_base = ''.join(ch if (ch.isalnum() or ch in {'-', '_'}) else '_' for ch in base_name).strip('_')
        if not safe_base:
            safe_base = f'rule_{rule.id}'
        safe_base = safe_base[:80]

        if format_type == 'json':
            payload = {
                'rule_id': rule.id,
                'rule_name': rule.name or '',
                'message_id': message.id,
                'role': message.role or 'assistant',
                'created_at': format_datetime_value(message.created_at),
                'generated_at': generated_at,
                'content': body
            }
            content = json.dumps(payload, ensure_ascii=False, indent=2)
            filename = f"{safe_base}_review_{message.id}.json"
        elif format_type == 'txt':
            content = (
                f"Rule ID: {rule.id}\n"
                f"Rule Name: {rule.name or ''}\n"
                f"Review Message ID: {message.id}\n"
                f"Role: {message.role or 'assistant'}\n"
                f"Created At: {format_datetime_value(message.created_at)}\n"
                f"Generated At: {generated_at}\n\n"
                f"{body}\n"
            )
            filename = f"{safe_base}_review_{message.id}.txt"
        else:
            content = (
                f"# Rule Revision Draft\n\n"
                f"- Rule ID: `{rule.id}`\n"
                f"- Rule Name: `{rule.name or ''}`\n"
                f"- Review Message ID: `{message.id}`\n"
                f"- Created At: `{format_datetime_value(message.created_at)}`\n\n"
                f"{body}\n"
            )
            filename = f"{safe_base}_review_{message.id}.md"

        return serialize_download_response(content, format_type, filename)
    except Exception:
        app.logger.exception('下载规则审核消息异常')
        return jsonify({'success': False, 'message': '下载失败'}), 500


@app.route('/api/messages/<int:history_id>/download', methods=['GET'])
@login_required
def download_message(history_id):
    try:
        format_type = (request.args.get('format') or 'md').strip().lower()
        if format_type not in {'md', 'txt', 'json'}:
            return jsonify({'success': False, 'message': '不支持的导出格式'}), 400

        chat = ChatHistory.query.filter_by(id=history_id, user_id=current_user.id).first()
        if not chat:
            return jsonify({'success': False, 'message': '消息不存在'}), 404

        attachments = resolve_chat_attachments(chat)
        attachment_names = [item.original_name for item in attachments]
        generated_at = datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S')

        if format_type == 'json':
            payload = {
                'id': chat.id,
                'conversation_id': chat.conversation_id or '',
                'model': chat.model or '',
                'created_at': format_datetime_value(chat.created_at),
                'generated_at': generated_at,
                'question': chat.question,
                'answer': chat.answer,
                'attachments': attachment_names
            }
            content = json.dumps(payload, ensure_ascii=False, indent=2)
            filename = f"reply_{chat.id}.json"
        elif format_type == 'txt':
            content = (
                f"Message ID: {chat.id}\n"
                f"Conversation: {chat.conversation_id or ''}\n"
                f"Model: {chat.model or ''}\n"
                f"Created At: {format_datetime_value(chat.created_at)}\n"
                f"Attachments: {', '.join(attachment_names) if attachment_names else 'None'}\n"
                f"Generated At: {generated_at}\n\n"
                "Question:\n"
                f"{chat.question}\n\n"
                "Answer:\n"
                f"{chat.answer}\n"
            )
            filename = f"reply_{chat.id}.txt"
        else:
            content = (
                f"# Reply #{chat.id}\n\n"
                f"- Conversation: `{chat.conversation_id or ''}`\n"
                f"- Model: `{chat.model or ''}`\n"
                f"- Created At: `{format_datetime_value(chat.created_at)}`\n"
                f"- Attachments: {', '.join(attachment_names) if attachment_names else 'None'}\n\n"
                "## Question\n\n"
                f"{chat.question}\n\n"
                "## Answer\n\n"
                f"{chat.answer}\n"
            )
            filename = f"reply_{chat.id}.md"

        return serialize_download_response(content, format_type, filename)
    except Exception:
        app.logger.exception('下载单条消息异常')
        return jsonify({'success': False, 'message': '下载失败'}), 500

@app.route('/api/chat', methods=['POST'])
@login_required
def chat():
    try:
        limit_response = enforce_rate_limit(
            'chat',
            f"{current_user.id}:{get_client_ip()}",
            limit=20,
            window_seconds=60
        )
        if limit_response:
            return limit_response

        data = request.get_json(silent=True) or {}
        if not data:
            return jsonify({'success': False, 'message': '无效的请求数据'})
        
        question = data.get('question', '').strip()
        prompt = (data.get('prompt') or '').strip()
        model = (data.get('model') or app.config['DEFAULT_MODEL']).strip()
        conversation_id = (data.get('conversation_id') or '').strip()
        attachment_ids = parse_attachment_ids(data.get('attachment_ids'))
        
        if not question:
            return jsonify({'success': False, 'message': '问题不能为空'})
        if not model:
            return jsonify({'success': False, 'message': '请先选择模型'})
        
        if len(question) > app.config['MAX_QUESTION_CHARS']:
            return jsonify({'success': False, 'message': f"问题过长，最多 {app.config['MAX_QUESTION_CHARS']} 字符"})
        if not prompt:
            prompt = question
        if len(prompt) > app.config['MAX_PROMPT_CHARS']:
            return jsonify({'success': False, 'message': f"请求上下文过长，最多 {app.config['MAX_PROMPT_CHARS']} 字符"})
        if len(attachment_ids) > app.config['MAX_ATTACHMENTS_PER_MESSAGE']:
            return jsonify({'success': False, 'message': f"单次最多关联 {app.config['MAX_ATTACHMENTS_PER_MESSAGE']} 个附件"})

        if conversation_id and (len(conversation_id) > 64 or '/' in conversation_id):
            return jsonify({'success': False, 'message': '会话标识无效'})
        if not conversation_id:
            conversation_id = secrets.token_hex(12)

        conversation_title = ''
        if conversation_id.startswith('legacy-'):
            try:
                legacy_id = int(conversation_id.split('-', 1)[1])
                legacy_chat = ChatHistory.query.filter_by(id=legacy_id, user_id=current_user.id).first()
                if legacy_chat:
                    conversation_title = (legacy_chat.conversation_title or '').strip() or (legacy_chat.question or '')[:60]
            except (IndexError, ValueError):
                return jsonify({'success': False, 'message': '会话标识无效'})
        else:
            existing_chat = (
                ChatHistory.query
                .filter_by(user_id=current_user.id, conversation_id=conversation_id)
                .order_by(ChatHistory.created_at.desc())
                .first()
            )
            if existing_chat:
                conversation_title = (existing_chat.conversation_title or '').strip() or (existing_chat.question or '')[:60]

        if not conversation_title:
            conversation_title = question[:60]

        attachments = []
        if attachment_ids:
            attachments = (
                Attachment.query
                .filter(Attachment.user_id == current_user.id, Attachment.id.in_(attachment_ids))
                .all()
            )
            found_ids = {item.id for item in attachments}
            if any(attachment_id not in found_ids for attachment_id in attachment_ids):
                return jsonify({'success': False, 'message': '附件不存在或无访问权限'})

            prompt_budget = app.config['MAX_PROMPT_CHARS'] - len(prompt) - 256
            if prompt_budget <= 0:
                return jsonify({'success': False, 'message': '问题上下文已接近上限，无法附加附件内容'})
            snippets = attachment_prompt_snippets(attachments, prompt_budget)
            if snippets:
                attachment_block = "\n\n".join(snippets)
                prompt = (
                    f"{prompt}\n\n"
                    "以下是用户上传附件的文本摘录（可能不完整，请结合问题判断）：\n"
                    f"{attachment_block}"
                )

        result = query_ollama(prompt, model)
        
        if 'error' in result:
            return jsonify({'success': False, 'message': result['error']})
        
        # 渲染Markdown
        raw_answer = result['response']
        html_answer = render_markdown(raw_answer)
        html_question = render_markdown(question)
        
        chat = ChatHistory(
            user_id=current_user.id,
            conversation_id=conversation_id,
            conversation_title=conversation_title,
            question=question,
            question_html=html_question,
            answer=raw_answer,
            answer_html=html_answer,
            attachment_ids=json.dumps(attachment_ids, ensure_ascii=False),
            model=model,
            tokens_used=result.get('eval_count', 0)
        )
        db.session.add(chat)
        for attachment in attachments:
            if not attachment.conversation_id and not conversation_id.startswith('legacy-'):
                attachment.conversation_id = conversation_id
        db.session.commit()
        
        return jsonify({
            'success': True,
            'question': question,
            'question_html': html_question,
            'answer': raw_answer,
            'answer_html': html_answer,
            'history_id': chat.id,
            'conversation_id': conversation_id,
            'conversation_title': conversation_title,
            'attachments': [item.to_dict() for item in attachments],
            'tokens_used': result.get('eval_count', 0)
        })
    except Exception:
        app.logger.exception('聊天处理异常')
        return jsonify({'success': False, 'message': '处理失败，请稍后重试'}), 500

@app.route('/api/history', methods=['GET'])
@login_required
def get_history():
    try:
        search = request.args.get('search', '').strip()
        summary = request.args.get('summary', '0').strip() in {'1', 'true', 'yes'}
        try:
            limit = int(request.args.get('limit', '300'))
        except ValueError:
            limit = 300
        limit = max(1, min(limit, 1000))
        
        query = ChatHistory.query.filter_by(user_id=current_user.id)
        
        if search:
            query = query.filter(
                (ChatHistory.question.ilike(f'%{search}%')) |
                (ChatHistory.answer.ilike(f'%{search}%'))
            )

        fetch_limit = limit
        if summary:
            fetch_limit = min(max(limit * 20, limit), 5000)

        history = query.order_by(ChatHistory.created_at.desc()).limit(fetch_limit).all()

        def conversation_key(item):
            conv = (item.conversation_id or '').strip()
            return conv if conv else f"legacy-{item.id}"

        history_list = []
        if summary:
            seen = set()
            for item in history:
                try:
                    conv_id = conversation_key(item)
                    if conv_id in seen:
                        continue
                    seen.add(conv_id)
                    history_list.append({
                        'conversation_id': conv_id,
                        'title': (item.conversation_title or '').strip() or (item.question or '')[:60],
                        'question': item.question,
                        'created_at': format_datetime_value(item.created_at),
                        'model': item.model or ''
                    })
                    if len(history_list) >= limit:
                        break
                except Exception:
                    app.logger.exception('历史摘要序列化失败，跳过 id=%s', getattr(item, 'id', None))
                    continue
        else:
            pending_render_updates = False
            for item in history:
                try:
                    history_data = item.to_dict()
                    # 如果没有HTML版本，就染一个
                    if not history_data.get('question_html'):
                        rendered = render_markdown(item.question)
                        history_data['question_html'] = rendered
                        item.question_html = rendered
                        pending_render_updates = True
                    if not history_data.get('answer_html'):
                        rendered = render_markdown(item.answer)
                        history_data['answer_html'] = rendered
                        item.answer_html = rendered
                        pending_render_updates = True
                    history_data['conversation_id'] = conversation_key(item)
                    history_list.append(history_data)
                except Exception:
                    app.logger.exception('历史记录序列化失败，跳过 id=%s', getattr(item, 'id', None))
                    continue

            if pending_render_updates:
                db.session.commit()
        
        return jsonify({
            'success': True,
            'history': history_list
        })
    except Exception:
        app.logger.exception('获取历史记录异常')
        return jsonify({'success': False, 'message': '获取历史记录失败'}), 500

@app.route('/api/history/<int:history_id>', methods=['GET', 'DELETE'])
@login_required
def delete_history(history_id):
    try:
        chat = ChatHistory.query.filter_by(id=history_id, user_id=current_user.id).first()
        if not chat:
            return jsonify({'success': False, 'message': '记录不存在'}), 404

        if request.method == 'GET':
            history_data = chat.to_dict()
            pending_render_updates = False
            if not history_data.get('question_html'):
                rendered = render_markdown(chat.question)
                history_data['question_html'] = rendered
                chat.question_html = rendered
                pending_render_updates = True
            if not history_data.get('answer_html'):
                rendered = render_markdown(chat.answer)
                history_data['answer_html'] = rendered
                chat.answer_html = rendered
                pending_render_updates = True
            if pending_render_updates:
                db.session.commit()
            history_data['conversation_id'] = (chat.conversation_id or '').strip() or f"legacy-{chat.id}"
            history_data['conversation_title'] = (chat.conversation_title or '').strip() or (chat.question or '')[:60]
            return jsonify({'success': True, 'history': history_data})

        db.session.delete(chat)
        db.session.commit()
        return jsonify({'success': True})
    except Exception:
        app.logger.exception('历史记录操作异常')
        if request.method == 'GET':
            return jsonify({'success': False, 'message': '获取记录失败'}), 500
        return jsonify({'success': False, 'message': '删除失败'}), 500


@app.route('/api/conversations/<conversation_id>', methods=['GET', 'DELETE'])
@login_required
def conversation_detail(conversation_id):
    try:
        conversation_id = (conversation_id or '').strip()
        if not conversation_id:
            return jsonify({'success': False, 'message': '会话标识不能为空'}), 400

        query = ChatHistory.query.filter_by(user_id=current_user.id)
        if conversation_id.startswith('legacy-'):
            try:
                legacy_id = int(conversation_id.split('-', 1)[1])
            except (IndexError, ValueError):
                return jsonify({'success': False, 'message': '会话标识无效'}), 400
            query = query.filter_by(id=legacy_id)
        else:
            query = query.filter_by(conversation_id=conversation_id)

        chats = query.order_by(ChatHistory.created_at.asc()).all()
        if not chats:
            return jsonify({'success': False, 'message': '会话不存在'}), 404

        if request.method == 'DELETE':
            for chat in chats:
                db.session.delete(chat)
            db.session.commit()
            return jsonify({'success': True})

        history_list = []
        pending_render_updates = False
        effective_title = ''
        for chat in chats:
            data = chat.to_dict()
            if not data.get('question_html'):
                rendered = render_markdown(chat.question)
                data['question_html'] = rendered
                chat.question_html = rendered
                pending_render_updates = True
            if not data.get('answer_html'):
                rendered = render_markdown(chat.answer)
                data['answer_html'] = rendered
                chat.answer_html = rendered
                pending_render_updates = True
            data['conversation_id'] = (chat.conversation_id or '').strip() or f"legacy-{chat.id}"
            data['conversation_title'] = (chat.conversation_title or '').strip() or (chat.question or '')[:60]
            if not effective_title and data['conversation_title']:
                effective_title = data['conversation_title']
            history_list.append(data)

        if pending_render_updates:
            db.session.commit()

        return jsonify({
            'success': True,
            'conversation_id': conversation_id,
            'conversation_title': effective_title,
            'history': history_list
        })
    except Exception:
        app.logger.exception('会话记录操作异常')
        if request.method == 'DELETE':
            return jsonify({'success': False, 'message': '删除会话失败'}), 500
        return jsonify({'success': False, 'message': '获取会话失败'}), 500


@app.route('/api/conversations/<conversation_id>/export', methods=['GET'])
@login_required
def export_conversation(conversation_id):
    try:
        format_type = (request.args.get('format') or 'md').strip().lower()
        if format_type not in {'md', 'txt', 'json'}:
            return jsonify({'success': False, 'message': '不支持的导出格式'}), 400

        conversation_id = (conversation_id or '').strip()
        if not conversation_id:
            return jsonify({'success': False, 'message': '会话标识不能为空'}), 400

        query = ChatHistory.query.filter_by(user_id=current_user.id)
        if conversation_id.startswith('legacy-'):
            try:
                legacy_id = int(conversation_id.split('-', 1)[1])
            except (IndexError, ValueError):
                return jsonify({'success': False, 'message': '会话标识无效'}), 400
            query = query.filter_by(id=legacy_id)
        else:
            query = query.filter_by(conversation_id=conversation_id)

        chats = query.order_by(ChatHistory.created_at.asc()).all()
        if not chats:
            return jsonify({'success': False, 'message': '会话不存在'}), 404

        conversation_title = ((chats[0].conversation_title or '').strip() or f"conversation-{conversation_id}")[:80]
        generated_at = datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S')

        if format_type == 'json':
            payload = {
                'conversation_id': conversation_id,
                'title': conversation_title,
                'generated_at': generated_at,
                'messages': []
            }
            for chat in chats:
                attachments = resolve_chat_attachments(chat)
                payload['messages'].append({
                    'id': chat.id,
                    'created_at': format_datetime_value(chat.created_at),
                    'model': chat.model or '',
                    'question': chat.question,
                    'answer': chat.answer,
                    'attachments': [item.original_name for item in attachments]
                })
            content = json.dumps(payload, ensure_ascii=False, indent=2)
            filename = f"conversation_{conversation_id}.json"
        elif format_type == 'txt':
            parts = [
                f"Conversation: {conversation_id}",
                f"Title: {conversation_title}",
                f"Generated At: {generated_at}",
                ''
            ]
            for index, chat in enumerate(chats, start=1):
                attachments = resolve_chat_attachments(chat)
                parts.extend([
                    f"=== Round {index} / Message {chat.id} ===",
                    f"Created At: {format_datetime_value(chat.created_at)}",
                    f"Model: {chat.model or ''}",
                    f"Attachments: {', '.join([item.original_name for item in attachments]) if attachments else 'None'}",
                    '',
                    "Question:",
                    chat.question,
                    '',
                    "Answer:",
                    chat.answer,
                    ''
                ])
            content = '\n'.join(parts)
            filename = f"conversation_{conversation_id}.txt"
        else:
            parts = [
                f"# {conversation_title}",
                '',
                f"- Conversation ID: `{conversation_id}`",
                f"- Generated At: `{generated_at}`",
                ''
            ]
            for index, chat in enumerate(chats, start=1):
                attachments = resolve_chat_attachments(chat)
                attachment_text = ', '.join([item.original_name for item in attachments]) if attachments else 'None'
                parts.extend([
                    f"## Round {index} (Message {chat.id})",
                    '',
                    f"- Time: `{format_datetime_value(chat.created_at)}`",
                    f"- Model: `{chat.model or ''}`",
                    f"- Attachments: {attachment_text}",
                    '',
                    "### Question",
                    '',
                    chat.question,
                    '',
                    "### Answer",
                    '',
                    chat.answer,
                    ''
                ])
            content = '\n'.join(parts)
            filename = f"conversation_{conversation_id}.md"

        return serialize_download_response(content, format_type, filename)
    except Exception:
        app.logger.exception('导出会话异常')
        return jsonify({'success': False, 'message': '导出失败'}), 500


@app.route('/api/conversations/<conversation_id>/title', methods=['PATCH'])
@login_required
def update_conversation_title(conversation_id):
    try:
        conversation_id = (conversation_id or '').strip()
        if not conversation_id:
            return jsonify({'success': False, 'message': '会话标识不能为空'}), 400

        data = request.get_json(silent=True) or {}
        title = (data.get('title') or '').strip()
        if not title:
            return jsonify({'success': False, 'message': '标题不能为空'}), 400
        if len(title) > 120:
            return jsonify({'success': False, 'message': '标题长度不能超过120字符'}), 400

        query = ChatHistory.query.filter_by(user_id=current_user.id)
        if conversation_id.startswith('legacy-'):
            try:
                legacy_id = int(conversation_id.split('-', 1)[1])
            except (IndexError, ValueError):
                return jsonify({'success': False, 'message': '会话标识无效'}), 400
            query = query.filter_by(id=legacy_id)
        else:
            query = query.filter_by(conversation_id=conversation_id)

        chats = query.all()
        if not chats:
            return jsonify({'success': False, 'message': '会话不存在'}), 404

        for chat in chats:
            chat.conversation_title = title
        db.session.commit()

        return jsonify({
            'success': True,
            'conversation_id': conversation_id,
            'conversation_title': title
        })
    except Exception:
        app.logger.exception('更新会话标题异常')
        return jsonify({'success': False, 'message': '更新会话标题失败'}), 500

@app.route('/api/clear_history', methods=['DELETE'])
@login_required
def clear_history():
    try:
        ChatHistory.query.filter_by(user_id=current_user.id).delete()
        db.session.commit()
        return jsonify({'success': True})
    except Exception:
        app.logger.exception('清空历史记录异常')
        return jsonify({'success': False, 'message': '清空失败'}), 500

@app.route('/api/models', methods=['GET'])
def get_models():
    """获取可用的模型列表"""
    try:
        last_status = None
        saw_success_response = False
        for base_url in get_ollama_base_urls():
            try:
                response = safe_requests_get(f"{base_url}/api/tags", timeout=5)
                if response.status_code != 200:
                    last_status = response.status_code
                    continue
                saw_success_response = True
                try:
                    data = response.json()
                except ValueError:
                    app.logger.warning('Ollama /api/tags 返回了非 JSON 响应: %s', base_url)
                    continue

                model_names = extract_model_names(data)
                if model_names:
                    print(f"✓ 获取到模型列表: {model_names}")
                    return jsonify({
                        'success': True,
                        'models': model_names
                    })
            except requests.exceptions.RequestException as exc:
                app.logger.warning('Ollama /api/tags 请求失败: %s (%s)', base_url, exc)
                continue
        if last_status:
            app.logger.warning('Ollama /api/tags 返回非200: %s', last_status)
        if saw_success_response:
            app.logger.info('Ollama /api/tags 可访问，但当前无可用模型')
    except Exception:
        app.logger.exception('获取模型列表失败')
    
    # 如果获取失败，回退到配置的默认模型（若有）
    fallback_models = [app.config['DEFAULT_MODEL']] if app.config['DEFAULT_MODEL'] else []
    return jsonify({
        'success': True,
        'models': fallback_models
    })

@app.route('/health')
def health():
    return jsonify({
        'status': 'ok',
        'authenticated': current_user.is_authenticated,
        'user': current_user.username if current_user.is_authenticated else None
    })


@app.route('/favicon.ico')
def favicon():
    return send_from_directory(app.static_folder, 'favicon.svg', mimetype='image/svg+xml')

# ============ Gunicorn入口 ============
def create_app():
    return app

# ============ 主程序 ============
if __name__ == '__main__':
    print("\n" + "="*50)
    print("🚀 Ollama WebUI 启动")
    print("="*50)
    print(f"📁 数据库: {app.config['SQLALCHEMY_DATABASE_URI']}")
    print(f"🌐 Ollama: {app.config['OLLAMA_BASE_URL']}")
    print(f"🤖 默认模型: {app.config['DEFAULT_MODEL'] or '（未设置）'}")
    print(f"🔐 调试模式: {env_flag('FLASK_DEBUG', False)}")
    print("="*50 + "\n")
    app.run(
        host=os.environ.get('FLASK_HOST', '127.0.0.1'),
        port=int(os.environ.get('PORT', '5001')),
        debug=env_flag('FLASK_DEBUG', False)
    )
